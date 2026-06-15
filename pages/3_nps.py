import io

import matplotlib.pyplot as plt
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt

from auth import init_auth, login_form, require_login
from theme import apply_theme, render_sidebar

COMPANY_NAME = "Pizzeria Insights"

apply_theme()
init_auth()

if not require_login():
    login_form(COMPANY_NAME)
    st.stop()

render_sidebar(COMPANY_NAME)


def load_data(uploaded_file):
    if uploaded_file.name.endswith(".csv"):
        return pd.read_csv(uploaded_file, encoding="utf-8")

    return pd.read_excel(uploaded_file)


def calculate_nps(df, column_name):
    nps_data = pd.to_numeric(df[column_name], errors="coerce").dropna()
    total_responses = len(nps_data)

    if total_responses == 0:
        return None

    promoters = int(nps_data[nps_data >= 9].count())
    passives = int(nps_data[(nps_data >= 7) & (nps_data <= 8)].count())
    detractors = int(nps_data[nps_data <= 6].count())

    promoter_percent = (promoters / total_responses) * 100
    passive_percent = (passives / total_responses) * 100
    detractor_percent = (detractors / total_responses) * 100
    nps_score = promoter_percent - detractor_percent

    return {
        "nps_score": nps_score,
        "promoters": promoters,
        "passives": passives,
        "detractors": detractors,
        "promoter_percent": promoter_percent,
        "passive_percent": passive_percent,
        "detractor_percent": detractor_percent,
        "total_responses": total_responses
    }


def create_results_df(result):
    return pd.DataFrame({
        "الفئة": ["المروجون", "المحايدون", "المنتقدون"],
        "العدد": [
            result["promoters"],
            result["passives"],
            result["detractors"]
        ],
        "النسبة المئوية": [
            f"{result['promoter_percent']:.2f}%",
            f"{result['passive_percent']:.2f}%",
            f"{result['detractor_percent']:.2f}%"
        ]
    })


def create_pptx_with_nps_results(result, nps_column):
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"NPS للعمود: {nps_column}"

    blank_layout = prs.slide_layouts[5]

    result_slide = prs.slides.add_slide(blank_layout)
    result_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    ).text_frame.text = "نتائج NPS"

    tf = result_slide.shapes.add_textbox(
        Inches(1), Inches(1.2), Inches(8), Inches(4)
    ).text_frame

    p = tf.paragraphs[0]
    p.text = f"صافي نقاط الترويج (NPS): {result['nps_score']:.2f}"
    p.font.bold = True
    p.font.size = Pt(18)

    tf.add_paragraph().text = f"النسبة المئوية للمروجين: {result['promoter_percent']:.2f}%"
    tf.add_paragraph().text = f"النسبة المئوية للمحايدين: {result['passive_percent']:.2f}%"
    tf.add_paragraph().text = f"النسبة المئوية للمنتقدين: {result['detractor_percent']:.2f}%"
    tf.add_paragraph().text = f"حجم العينة: {result['total_responses']}"

    chart_slide = prs.slides.add_slide(blank_layout)
    chart_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    ).text_frame.text = "توزيع فئات NPS"

    chart_data = pd.DataFrame({
        "الفئة": ["المروجون", "المحايدون", "المنتقدون"],
        "النسبة المئوية": [
            result["promoter_percent"],
            result["passive_percent"],
            result["detractor_percent"]
        ]
    })

    fig, ax = plt.subplots(figsize=(8, 5))
    ax.bar(
        chart_data["الفئة"],
        chart_data["النسبة المئوية"],
        color=["green", "orange", "red"]
    )
    ax.set_title("توزيع فئات NPS")
    ax.set_ylabel("النسبة المئوية")
    plt.tight_layout()

    chart_image = io.BytesIO()
    fig.savefig(chart_image, format="png", bbox_inches="tight")
    chart_image.seek(0)
    plt.close(fig)

    chart_slide.shapes.add_picture(chart_image, Inches(1), Inches(1.5), width=Inches(8))

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return output


st.title("📊 NPS Analysis")

uploaded_file = st.file_uploader("⬆️ Upload Excel أو CSV", type=["xlsx", "xls", "csv"])

if uploaded_file is not None:
    try:
        df = load_data(uploaded_file)
        df.columns = df.columns.str.strip()

        nps_column = st.selectbox("اختر العمود المراد تحليله:", df.columns.tolist())

        if nps_column:
            st.subheader(f"تحليل العمود: {nps_column}")

            result = calculate_nps(df, nps_column)

            if result is None:
                st.warning("لا توجد بيانات رقمية كافية في العمود المحدد.")
            else:
                st.markdown(f"**حجم العينة:** {result['total_responses']}")

                results_df = create_results_df(result)

                st.dataframe(results_df, use_container_width=True)

                chart_counts_df = pd.DataFrame({
                    "الفئة": ["المروجون", "المحايدون", "المنتقدون"],
                    "العدد": [
                        result["promoters"],
                        result["passives"],
                        result["detractors"]
                    ]
                })

                st.bar_chart(chart_counts_df.set_index("الفئة")["العدد"])

                st.metric(
                    label="صافي نقاط الترويج (NPS)",
                    value=f"{result['nps_score']:.2f}"
                )

                if result["nps_score"] >= 50:
                    st.success("النتيجة ممتازة! لديك قاعدة عملاء مخلصين جدًا.")
                elif result["nps_score"] >= 0:
                    st.info("النتيجة جيدة. لديك عدد من العملاء الراضين.")
                else:
                    st.warning("النتيجة تحتاج إلى تحسين.")

                st.subheader("تصدير النتائج")

                excel_output = io.BytesIO()
                with pd.ExcelWriter(excel_output, engine="xlsxwriter") as writer:
                    results_df.to_excel(writer, sheet_name="NPS_Analysis", index=False)

                st.download_button(
                    label="تصدير إلى Excel",
                    data=excel_output.getvalue(),
                    file_name="NPS_Analysis.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

                pptx_output = create_pptx_with_nps_results(result, nps_column)

                st.download_button(
                    label="تصدير إلى PowerPoint",
                    data=pptx_output.getvalue(),
                    file_name="NPS_Presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

    except Exception as e:
        st.error(f"حدث خطأ أثناء معالجة الملف: {e}")
