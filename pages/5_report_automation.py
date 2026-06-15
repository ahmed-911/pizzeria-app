import io
from datetime import date
from pathlib import Path

import pandas as pd
import streamlit as st
from pptx import Presentation

from auth import init_auth, login_form, require_login
from theme import apply_theme, render_sidebar

COMPANY_NAME = "Pizzeria Insights"

PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = PROJECT_ROOT / "templates" / "weekly_report_template.pptx"
PREVIEW_PATH = PROJECT_ROOT / "templates" / "weekly_report_preview.png"

apply_theme()
init_auth()

if not require_login():
    login_form(COMPANY_NAME)
    st.stop()

render_sidebar(COMPANY_NAME)


def load_data(uploaded_file):
    if uploaded_file.name.endswith(".csv"):
        return pd.read_csv(uploaded_file, encoding="utf-8")

    return pd.read_excel(uploaded_file)


def calculate_csat(df, column_name):
    data = pd.to_numeric(df[column_name], errors="coerce").dropna()
    total = len(data)

    if total == 0:
        return 0

    satisfied = (data >= 4).sum()
    return round((satisfied / total) * 100, 1)


def calculate_ces(df, column_name):
    data = pd.to_numeric(df[column_name], errors="coerce").dropna()
    total = len(data)

    if total == 0:
        return 0

    satisfied = (data >= 5).sum()
    return round((satisfied / total) * 100, 1)


def calculate_nps(df, column_name):
    data = pd.to_numeric(df[column_name], errors="coerce").dropna()
    total = len(data)

    if total == 0:
        return 0

    promoters = (data >= 9).sum()
    detractors = (data <= 6).sum()

    promoter_percent = (promoters / total) * 100
    detractor_percent = (detractors / total) * 100

    return round(promoter_percent - detractor_percent, 1)


def replace_placeholders(template_path, values):
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)

                for placeholder, value in values.items():
                    full_text = full_text.replace(placeholder, str(value))

                if paragraph.runs:
                    paragraph.runs[0].text = full_text
                    for run in paragraph.runs[1:]:
                        run.text = ""

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return output


def build_placeholder_reference():
    return pd.DataFrame({
        "Placeholder": [
            "{{CSAT_RECEPTION}}",
            "{{CSAT_DELIVERY}}",
            "{{CSAT_FOOD}}",
            "{{CSAT_OVERALL}}",
            "{{CES_SCORE}}",
            "{{NPS_SCORE}}",
            "{{TOTAL_RESPONSES}}",
            "{{REPORT_DATE}}",
        ],
        "Description": [
            "Reception satisfaction score",
            "Delivery satisfaction score",
            "Food quality satisfaction score",
            "Overall satisfaction score",
            "Customer effort score",
            "Net promoter score",
            "Total survey responses",
            "Report date",
        ]
    })


st.title("📄 Report Automation")

if not TEMPLATE_PATH.exists():
    st.error("PowerPoint template not found.")
    st.write("Looking for template here:")
    st.code(str(TEMPLATE_PATH))
    st.info("Please add your template here: templates/weekly_report_template.pptx")
    st.stop()

with st.expander("Template Preview & Placeholder Reference", expanded=True):
    if PREVIEW_PATH.exists():
        st.image(
            str(PREVIEW_PATH),
            caption="Report Template Preview",
            use_container_width=True
        )
    else:
        st.info("Optional: add a preview image here: templates/weekly_report_preview.png")

    st.dataframe(build_placeholder_reference(), use_container_width=True)

uploaded_file = st.file_uploader(
    "Upload Weekly Excel or CSV File",
    type=["xlsx", "xls", "csv"]
)

if uploaded_file:
    try:
        df = load_data(uploaded_file)
        df.columns = df.columns.astype(str).str.strip()

        st.subheader("Select Columns")

        csat_reception_col = st.selectbox("CSAT Reception Column", df.columns)
        csat_delivery_col = st.selectbox("CSAT Delivery Column", df.columns)
        csat_food_col = st.selectbox("CSAT Food Column", df.columns)
        csat_overall_col = st.selectbox("CSAT Overall Column", df.columns)

        ces_column = st.selectbox("CES Column", df.columns)
        nps_column = st.selectbox("NPS Column", df.columns)

        report_date = st.date_input("Report Date", value=date.today())

        if st.button("Generate Automated Report", use_container_width=True):
            values = {
                "{{CSAT_RECEPTION}}": f"{calculate_csat(df, csat_reception_col)}%",
                "{{CSAT_DELIVERY}}": f"{calculate_csat(df, csat_delivery_col)}%",
                "{{CSAT_FOOD}}": f"{calculate_csat(df, csat_food_col)}%",
                "{{CSAT_OVERALL}}": f"{calculate_csat(df, csat_overall_col)}%",
                "{{CES_SCORE}}": f"{calculate_ces(df, ces_column)}%",
                "{{NPS_SCORE}}": calculate_nps(df, nps_column),
                "{{TOTAL_RESPONSES}}": len(df),
                "{{REPORT_DATE}}": report_date.strftime("%Y-%m-%d"),
            }

            report_output = replace_placeholders(TEMPLATE_PATH, values)

            st.success("Report generated successfully")

            st.download_button(
                label="Download PowerPoint Report",
                data=report_output.getvalue(),
                file_name="automated_weekly_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

    except Exception as e:
        st.error(f"Error while generating report: {e}")
