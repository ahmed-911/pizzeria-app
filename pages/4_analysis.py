import io

import matplotlib.pyplot as plt
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches

from auth import init_auth, login_form, require_login
from theme import apply_theme, render_sidebar

COMPANY_NAME = "Pizzeria Insights"

apply_theme()
init_auth()

if not require_login():
    login_form(COMPANY_NAME)
    st.stop()

render_sidebar(COMPANY_NAME)


def load_data(uploaded_file):
    if uploaded_file.name.endswith(".csv"):
        return pd.read_csv(uploaded_file, encoding="utf-8")

    return pd.read_excel(uploaded_file)


def create_ppt(title, df, fig):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]

    slide_chart = prs.slides.add_slide(slide_layout)
    tx_box = slide_chart.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    tx_box.text_frame.text = title

    img_stream = io.BytesIO()
    fig.savefig(img_stream, format="png", bbox_inches="tight")
    img_stream.seek(0)

    slide_chart.shapes.add_picture(
        img_stream,
        Inches(1),
        Inches(1.5),
        Inches(7),
        Inches(4)
    )

    slide_table = prs.slides.add_slide(slide_layout)
    tx_box2 = slide_table.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    tx_box2.text_frame.text = "Table"

    table_df = df.head(25)

    fig_table, ax_table = plt.subplots(figsize=(10, 6))
    ax_table.axis("off")

    tbl = ax_table.table(
        cellText=table_df.values,
        colLabels=table_df.columns,
        loc="center"
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(9)

    plt.tight_layout()

    img_table = io.BytesIO()
    fig_table.savefig(img_table, format="png", bbox_inches="tight")
    plt.close(fig_table)
    img_table.seek(0)

    slide_table.shapes.add_picture(
        img_table,
        Inches(0.5),
        Inches(1.5),
        width=Inches(9)
    )

    pptx_output = io.BytesIO()
    prs.save(pptx_output)
    pptx_output.seek(0)

    return pptx_output


def analyze_columns(df, selected_cols, chart_option):
    analysis_df = df[selected_cols].copy()
    analysis_df = analysis_df.fillna("Missing")

    count_col_name = "Count"
    while count_col_name in analysis_df.columns:
        count_col_name += "_1"

    counts = (
        analysis_df
        .groupby(selected_cols)
        .size()
        .reset_index(name=count_col_name)
    )

    if chart_option == "Percentages":
        counts["Percentage"] = (
            counts[count_col_name] / counts[count_col_name].sum() * 100
        ).round(2)

    st.subheader("Results")
    st.dataframe(counts, use_container_width=True)

    fig, ax = plt.subplots(figsize=(8, 4))

    labels = counts[selected_cols].astype(str).agg(" | ".join, axis=1)

    if chart_option == "Numbers":
        ax.bar(labels, counts[count_col_name])
        ax.set_ylabel("Count")
    else:
        ax.bar(labels, counts["Percentage"])
        ax.set_ylabel("Percentage")

    ax.set_title("Distribution")
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()

    st.pyplot(fig)

    excel_output = io.BytesIO()

    with pd.ExcelWriter(excel_output, engine="openpyxl") as writer:
        counts.to_excel(writer, sheet_name="Analysis", index=False)

    excel_output.seek(0)

    st.download_button(
        label="Download Excel",
        data=excel_output.getvalue(),
        file_name="analysis.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    ppt_output = create_ppt("Analysis Results", counts, fig)

    st.download_button(
        label="Download PowerPoint",
        data=ppt_output.getvalue(),
        file_name="analysis.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    plt.close(fig)


st.title("📊 Excel Column Analysis Tool")

uploaded_file = st.file_uploader(
    "Upload your Excel or CSV file",
    type=["xlsx", "xls", "csv"]
)

if uploaded_file:
    if uploaded_file.name.endswith(".csv"):
        df = load_data(uploaded_file)
    else:
        xls = pd.ExcelFile(uploaded_file)
        sheet_name = st.selectbox("Select a sheet", xls.sheet_names)
        df = pd.read_excel(uploaded_file, sheet_name=sheet_name)

    df.columns = df.columns.astype(str).str.strip()

    analysis_type = st.radio(
        "Select analysis type",
        ["One column", "Two columns", "Three columns", "Four columns", "Five columns"]
    )

    chart_option = st.radio("Chart Display", ["Numbers", "Percentages"])

    num_cols_map = {
        "One column": 1,
        "Two columns": 2,
        "Three columns": 3,
        "Four columns": 4,
        "Five columns": 5
    }

    num_cols = num_cols_map[analysis_type]

    if len(df.columns) < num_cols:
        st.warning("The file does not have enough columns for this analysis.")
    else:
        selected_cols = []
        available_columns = list(df.columns)

        for i in range(num_cols):
            col = st.selectbox(
                f"Select column {i + 1}",
                available_columns,
                key=f"analysis_col_{i}"
            )
            selected_cols.append(col)
            available_columns = [c for c in available_columns if c != col]

        if st.button("Run Analysis", use_container_width=True):
            analyze_columns(df, selected_cols, chart_option)
